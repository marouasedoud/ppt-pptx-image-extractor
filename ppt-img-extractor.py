import os
import comtypes.client
from pptx import Presentation
from docx import Document
from docx.shared import Inches
from PIL import Image
import io

# Function to convert .ppt to .pptx
def convert_ppt_to_pptx(input_ppt_file, output_pptx_file):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    ppt_file = os.path.abspath(input_ppt_file)
    pptx_file = os.path.abspath(output_pptx_file)
    deck = powerpoint.Presentations.Open(ppt_file)
    deck.SaveAs(pptx_file, 24)  # 24 represents the pptx file format
    deck.Close()
    powerpoint.Quit()
    return pptx_file

# List of .ppt files to process
ppt_files = ["File1.ppt", "File2.ppt","File3.ppt"]

# List to store resulting .pptx file paths
pptx_files = []

# Convert .ppt to .pptx
for ppt_file in ppt_files:
    pptx_output_file = f"{os.path.splitext(ppt_file)[0]}.pptx"
    pptx_files.append(convert_ppt_to_pptx(ppt_file, pptx_output_file))



# List to store extracted images
extracted_images = []

# Function to extract images from a PowerPoint file
def extract_images_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    extracted_images_from_file = []
    
    # Iterate through slides and extract images
    for slide in presentation.slides:
        slide_images = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text != "":  # Check if shape contains text
                continue
            if hasattr(shape, "image"):
                image_bytes = shape.image.blob
                image = Image.open(io.BytesIO(image_bytes))
                slide_images.append(image)
            # if len(slide_images) >= 2:
            #     break
        
        extracted_images_from_file.extend(slide_images[:2])
        # if len(extracted_images_from_file) >= 2:
        #     break
    
    return extracted_images_from_file

# Extract images from each PowerPoint file
for ppt_file in pptx_files:
    extracted_images.extend(extract_images_from_ppt(ppt_file))

# Create a new Word document
doc = Document()

# Add extracted images to the Word document
for i, image in enumerate(extracted_images):
    image_path = f"extracted_image_{i + 1}.png"
    image.save(image_path)
    doc.add_picture(image_path, width=Inches(2))

# Save the final document
doc.save("final_file.docx")

print("Extraction completed. Extracted images saved in final_file.docx.")

#########################################################################################################
# *************************************** WORKS PERFECTLY ***********************************************
#########################################################################################################
